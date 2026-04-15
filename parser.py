"""
文档解析器模块
支持多种文件格式的文本提取
"""
import logging
from pathlib import Path
from typing import Optional
import chardet

logger = logging.getLogger(__name__)


class DocumentParser:
    """文档解析器"""
    
    def __init__(self):
        pass
    
    def parse(self, file_path: Path) -> Optional[str]:
        """解析文件并提取文本内容"""
        try:
            suffix = file_path.suffix.lower()
            
            if suffix in ['.txt', '.md', '.rst', '.html', '.htm', '.xml', '.json', 
                         '.yaml', '.yml', '.log', '.ini', '.cfg', '.conf',
                         '.py', '.js', '.ts', '.java', '.cpp', '.c', '.h', '.cs', '.php', '.rb', '.go']:
                return self._parse_text(file_path)
            elif suffix in ['.docx', '.doc', '.docm', '.dotx', '.dotm']:
                return self._parse_word(file_path)
            elif suffix in ['.xlsx', '.xls', '.xlsm', '.et']:
                return self._parse_excel(file_path)
            elif suffix in ['.pptx', '.ppt', '.pptm', '.potx', '.potm', '.dps']:
                return self._parse_powerpoint(file_path)
            elif suffix == '.pdf':
                return self._parse_pdf(file_path)
            elif suffix == '.rtf':
                return self._parse_rtf(file_path)
            elif suffix in ['.odt', '.ods', '.odp']:
                return self._parse_openoffice(file_path)
            elif suffix in ['.epub', '.mobi', '.azw', '.azw3']:
                return self._parse_ebook(file_path)
            else:
                logger.debug(f"不支持的文件格式：{suffix}")
                return None
                
        except Exception as e:
            logger.error(f"解析文件失败 {file_path}: {e}")
            return None
    
    def _parse_text(self, file_path: Path) -> Optional[str]:
        """解析纯文本文件"""
        try:
            with open(file_path, 'rb') as f:
                raw_data = f.read()
            
            result = chardet.detect(raw_data)
            encoding = result['encoding'] if result['encoding'] else 'utf-8'
            
            try:
                text = raw_data.decode(encoding)
            except UnicodeDecodeError:
                text = raw_data.decode('utf-8', errors='ignore')
            
            return text
        except Exception as e:
            logger.error(f"解析文本文件失败 {file_path}: {e}")
            return None
    
    def _parse_word(self, file_path: Path) -> Optional[str]:
        """解析 Word 文档"""
        try:
            if file_path.suffix.lower() == '.docx':
                from docx import Document
                doc = Document(file_path)
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                return '\n'.join(paragraphs)
            elif file_path.suffix.lower() == '.doc':
                return self._parse_doc_legacy(file_path)
        except ImportError:
            logger.warning("python-docx 未安装，无法解析 Word 文档")
            return None
        except Exception as e:
            logger.error(f"解析 Word 文档失败 {file_path}: {e}")
            return None
        
        return None
    
    def _parse_doc_legacy(self, file_path: Path) -> Optional[str]:
        """解析旧版 .doc 文件（需要 antiword 或类似工具）"""
        try:
            import subprocess
            result = subprocess.run(
                ['antiword', str(file_path)],
                capture_output=True,
                text=True,
                timeout=30
            )
            if result.returncode == 0:
                return result.stdout
        except Exception as e:
            logger.warning(f"无法解析 .doc 文件 {file_path}: {e}")
        return None
    
    def _parse_excel(self, file_path: Path) -> Optional[str]:
        """解析 Excel 文件（包括 WPS 格式）"""
        try:
            suffix = file_path.suffix.lower()
            
            # 支持 .xlsx, .xls, 以及 WPS 的 .et 格式
            if suffix in ['.xlsx', '.et']:
                from openpyxl import load_workbook
                wb = load_workbook(file_path, read_only=True, data_only=True)
                
                texts = []
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        # 处理单元格内容，包括公式计算结果
                        row_values = []
                        for cell in row:
                            if cell is None:
                                row_values.append('')
                            elif isinstance(cell, (int, float)):
                                row_values.append(str(cell))
                            else:
                                row_values.append(str(cell))
                        
                        row_text = ' '.join(row_values)
                        if row_text.strip():
                            texts.append(row_text)
                
                result = '\n'.join(texts)
                if result.strip():
                    return result
                else:
                    logger.debug(f"Excel 文件无文本内容：{file_path}")
                    return ""
                    
            elif suffix == '.xls':
                return self._parse_xls_legacy(file_path)
            else:
                logger.warning(f"不支持的 Excel 格式：{suffix}")
                return None
                
        except ImportError as e:
            logger.error(f"缺少依赖库，无法解析 Excel 文档：{e}")
            return None
        except Exception as e:
            logger.error(f"解析 Excel 文档失败 {file_path}: {e}")
            return None
        
        return None
    
    def _parse_xls_legacy(self, file_path: Path) -> Optional[str]:
        """解析旧版 .xls 文件"""
        try:
            import xlrd
            wb = xlrd.open_workbook(file_path)
            
            texts = []
            for sheet in wb.sheets():
                for row in range(sheet.nrows):
                    row_values = sheet.row_values(row)
                    row_text = ' '.join(str(cell) for cell in row_values)
                    if row_text.strip():
                        texts.append(row_text)
                
                return '\n'.join(texts)
        except ImportError:
            logger.warning("xlrd 未安装，无法解析 .xls 文件")
            return None
        except Exception as e:
            logger.error(f"解析 .xls 文件失败 {file_path}: {e}")
            return None
        
        return None
    
    def _parse_powerpoint(self, file_path: Path) -> Optional[str]:
        """解析 PowerPoint 文件"""
        try:
            if file_path.suffix.lower() == '.pptx':
                from pptx import Presentation
                prs = Presentation(file_path)
                
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            texts.append(shape.text)
                
                return '\n'.join(texts)
            elif file_path.suffix.lower() == '.ppt':
                logger.warning(".ppt 格式需要额外工具支持")
                return None
        except ImportError:
            logger.warning("python-pptx 未安装，无法解析 PowerPoint 文档")
            return None
        except Exception as e:
            logger.error(f"解析 PowerPoint 文档失败 {file_path}: {e}")
            return None
        
        return None
    
    def _parse_ebook(self, file_path: Path) -> Optional[str]:
        """解析电子书文件（EPUB、MOBI、AZW3）"""
        try:
            suffix = file_path.suffix.lower()
            
            if suffix == '.epub':
                return self._parse_epub(file_path)
            elif suffix in ['.mobi', '.azw', '.azw3']:
                return self._parse_mobi(file_path)
            else:
                logger.warning(f"不支持的电子书格式：{suffix}")
                return None
        except Exception as e:
            logger.error(f"解析电子书失败 {file_path}: {e}")
            return None
    
    def _parse_epub(self, file_path: Path) -> Optional[str]:
        """解析 EPUB 电子书"""
        try:
            from ebooklib import epub
            import re
            
            book = epub.read_epub(file_path)
            texts = []
            
            # 提取所有章节内容
            for item in book.get_items():
                if item.get_type().name == 'ITEM_DOCUMENT':
                    # 提取 HTML 内容
                    content = item.get_content().decode('utf-8', errors='ignore')
                    
                    # 移除 HTML 标签
                    text = re.sub(r'<[^>]+>', ' ', content)
                    # 移除多余空白
                    text = re.sub(r'\s+', ' ', text).strip()
                    
                    if text:
                        texts.append(text)
            
            result = '\n'.join(texts[:100])  # 限制章节数量，避免太大
            return result if result else ""
        except ImportError:
            logger.error("缺少 ebooklib 库，无法解析 EPUB 文件")
            return None
        except Exception as e:
            logger.error(f"解析 EPUB 文件失败 {file_path}: {e}")
            return None
    
    def _parse_mobi(self, file_path: Path) -> Optional[str]:
        """解析 MOBI/AZW3 电子书"""
        try:
            # 尝试使用 mobi 库
            try:
                import mobi
                import tempfile
                import os
                
                # 解压到临时目录
                with tempfile.TemporaryDirectory() as temp_dir:
                    mobi_path, _ = mobi.extract(file_path, temp_dir)
                    
                    # 读取解压后的文本
                    texts = []
                    for root, dirs, files in os.walk(mobi_path):
                        for file in files:
                            if file.endswith('.html') or file.endswith('.xhtml'):
                                file_path_full = os.path.join(root, file)
                                with open(file_path_full, 'r', encoding='utf-8', errors='ignore') as f:
                                    content = f.read()
                                    # 移除 HTML 标签
                                    import re
                                    text = re.sub(r'<[^>]+>', ' ', content)
                                    text = re.sub(r'\s+', ' ', text).strip()
                                    if text:
                                        texts.append(text)
                    
                    result = '\n'.join(texts[:100])  # 限制内容量
                    return result if result else ""
            except ImportError:
                logger.warning("mobi 库未安装，尝试基础解析")
                # 基础解析：尝试直接读取文本内容
                try:
                    with open(file_path, 'rb') as f:
                        content = f.read()
                        # 尝试提取文本（简单方法，可能不完整）
                        text = content.decode('utf-8', errors='ignore')
                        # 移除控制字符
                        import re
                        text = re.sub(r'[\x00-\x08\x0B-\x0C\x0E-\x1F]', '', text)
                        return text[:50000]  # 限制长度
                except Exception:
                    return None
        except Exception as e:
            logger.error(f"解析 MOBI/AZW3 文件失败 {file_path}: {e}")
            return None
    
    def _parse_pdf(self, file_path: Path) -> Optional[str]:
        """解析 PDF 文件（带资源清理）"""
        try:
            import pdfplumber
            import threading
            
            logger.debug(f"开始解析 PDF: {file_path}")
            
            # 使用线程 + 超时控制
            result = [None]
            error = [None]
            
            def parse_pdf():
                try:
                    texts = []
                    # 使用上下文管理器确保资源清理
                    with pdfplumber.open(file_path) as pdf:
                        # 限制最多解析前 50 页，避免大文件卡住
                        page_count = min(len(pdf.pages), 50)
                        for i in range(page_count):
                            page = pdf.pages[i]
                            text = page.extract_text()
                            if text:
                                texts.append(text)
                        logger.debug(f"PDF 解析完成 {file_path}: {len(texts)} 页")
                    result[0] = '\n'.join(texts)
                except Exception as e:
                    error[0] = e
                    logger.error(f"PDF 解析错误 {file_path}: {e}")
            
            # 创建并启动线程
            thread = threading.Thread(target=parse_pdf, daemon=True)
            thread.start()
            thread.join(timeout=30)  # 30 秒超时
            
            if thread.is_alive():
                logger.warning(f"PDF 解析超时：{file_path}（超过 30 秒）")
                return f"[PDF 解析超时] 文件过大或内容复杂：{file_path.name}"
            
            if error[0]:
                logger.error(f"解析 PDF 文档失败 {file_path}: {error[0]}")
                return None
            
            result_text = result[0] if result[0] else ""
            logger.debug(f"PDF 解析成功 {file_path}: {len(result_text)} 字符")
            return result_text
            
        except ImportError:
            logger.warning("pdfplumber 未安装，无法解析 PDF 文档")
            return None
        except Exception as e:
            logger.error(f"解析 PDF 文档失败 {file_path}: {e}")
            return None
    
    def _parse_rtf(self, file_path: Path) -> Optional[str]:
        """解析 RTF 文件"""
        try:
            import striprtf.striprtf as striprtf
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                rtf_content = f.read()
            
            text = striprtf.rtf_to_text(rtf_content)
            return text
        except ImportError:
            logger.warning("striprtf 未安装，无法解析 RTF 文件")
            return None
        except Exception as e:
            logger.error(f"解析 RTF 文件失败 {file_path}: {e}")
            return None
    
    def _parse_openoffice(self, file_path: Path) -> Optional[str]:
        """解析 OpenOffice 文件"""
        try:
            import zipfile
            import xml.etree.ElementTree as ET
            
            with zipfile.ZipFile(file_path) as z:
                if 'content.xml' in z.namelist():
                    with z.open('content.xml') as f:
                        tree = ET.parse(f)
                        root = tree.getroot()
                        
                        texts = []
                        for elem in root.iter():
                            if elem.text and elem.text.strip():
                                texts.append(elem.text.strip())
                        
                        return ' '.join(texts)
        except Exception as e:
            logger.error(f"解析 OpenOffice 文件失败 {file_path}: {e}")
            return None


def extract_text(file_path: Path, max_size: int = 50 * 1024 * 1024) -> Optional[str]:
    """提取文件文本的便捷函数"""
    if not file_path.exists():
        logger.error(f"文件不存在：{file_path}")
        return None
    
    if file_path.stat().st_size > max_size:
        logger.warning(f"文件过大，跳过：{file_path}")
        return None
    
    parser = DocumentParser()
    return parser.parse(file_path)
